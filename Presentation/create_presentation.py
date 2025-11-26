"""
PowerPoint Presentation Generator for DIO Research Project
Creates a 15-minute presentation ready for export
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
import os

# Get the directory of this script
script_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(script_dir)

# Load optimization results from parent directory
with open(os.path.join(parent_dir, 'optimization_results.json'), 'r') as f:
    opt_results = json.load(f)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    return slide

def add_content_slide(prs, title, content_type='bullet'):
    """Add a content slide with title and body"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    return slide

def add_bullet_points(text_frame, points, level=0):
    """Add bullet points to a text frame"""
    text_frame.clear()
    for i, point in enumerate(points):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = point
        p.level = level
        p.font.size = Pt(18)
        p.space_after = Pt(12)

def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image"""
    slide = add_content_slide(prs, title)
    
    # Add image
    left = Inches(1.5)
    top = Inches(2)
    height = Inches(4.5)
    
    try:
        pic = slide.shapes.add_picture(image_path, left, top, height=height)
        
        # Add caption if provided
        if caption:
            left = Inches(1)
            top = Inches(6.5)
            width = Inches(8)
            height = Inches(0.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = caption
            p.font.size = Pt(14)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
    except:
        # If image not found, add placeholder text
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"[Image: {image_path}]"
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
add_title_slide(
    prs,
    "Dholes-Inspired Optimization (DIO)\nfor Feature Selection & Hyperparameter Tuning",
    "Cross-Domain Validation: From Medical Diagnostics to Computer Vision\n\nUSTO-MB | Computer Science | Data Science | Nov 2025"
)

# ============================================================================
# SLIDE 1A: DIO Authorship, Publication & Code
# ============================================================================
slide = add_content_slide(prs, "Who Created DIO? (Authorship & Code)")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "DIO created by:",
    "- Ali El Romeh (Centre for AI Research & Optimization, Torrens University Australia)",
    "- Seyedali Mirjalili (Lead researcher, Torrens University Australia)",
    "- Václav Šnel (VSB-Technical University of Ostrava, Czech Republic)",
    "",
    "Publication: Cluster Computing (Springer), 2025",
    "- DOI: 10.1007/s10586-025-05543-2",
    "- Received: Jan 27, 2025 | Accepted: May 12, 2025",
    "- High-tier peer-reviewed journal",
    "",
    "Code availability:",
    "- GitHub: github.com/AlyromehDholes-Inspired-Optimization-DIO",
    "- MathWorks File Exchange (MATLAB version)"
])

# ============================================================================
# SLIDE 1B: INTRO SLIDE (Names, University, Department, etc.)
# ============================================================================
slide = add_content_slide(prs, "Introduction & Team")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Names: Bellatreche Mohamed Amine, Iman Cherif Ghizlane",
    "University: USTO-MB (Université des Sciences et de la Technologie d'Oran Mohamed-Boudiaf)",
    "Department: Computer Science",
    "Specialty: Data Science",
    "Course: Statistics in Data Science",
    "Professor: Dr. Neggaz Nabil"
])

# ============================================================================
# SLIDE 1C: CONTENT PAGE (Table of Contents)
# ============================================================================
slide = add_content_slide(prs, "Contents")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Introduction & Motivation",
    "DIO Algorithm Overview",
    "Methodology: Nested Optimization",
    "Medical Classification Results (Breast Cancer)",
    "Extension to Computer Vision (CIFAR-10)",
    "Cross-Domain Analysis",
    "Statistical Validation",
    "Practical Implications",
    "Limitations & Future Work",
    "Conclusions",
    "References"
])

# ============================================================================
# SLIDE 2: AGENDA
# ============================================================================
slide = add_content_slide(prs, "Presentation Outline")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Introduction & Motivation",
    "DIO Algorithm Overview",
    "Methodology: Nested Optimization",
    "Medical Classification Results (Breast Cancer)",
    "Extension to Computer Vision (CIFAR-10)",
    "Cross-Domain Analysis",
    "Statistical Validation",
    "Practical Implications",
    "Limitations & Future Work",
    "Conclusions"
])

# ============================================================================
# SLIDE 3: PROBLEM STATEMENT
# ============================================================================
slide = add_content_slide(prs, "The Challenge We Faced")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Breast cancer: Leading cause of mortality in women worldwide",
    "Machine learning for diagnosis: High-dimensional data (30 features)",
    "The dilemma: Optimal feature selection + hyperparameter tuning simultaneously",
    "Traditional sequential optimization → Suboptimal, misses feature-parameter interactions",
    "Our approach: Simultaneous optimization using DIO (nature-inspired algorithm)"
])

# ============================================================================
# SLIDE 4: WHAT IS DIO?
# ============================================================================
slide = add_content_slide(prs, "Dholes-Inspired Optimization (DIO)")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Nature-inspired metaheuristic algorithm",
    "Based on pack hunting behavior of dholes (Asiatic wild dogs)",
    "Three cooperative hunting strategies:",
])

# ============================================================================
# SLIDE 4A: Dhole Animal & Behavior Mapping
# ============================================================================
slide = add_content_slide(prs, "Dholes: Animal Inspiration & Algorithm Mapping")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Dhole (Cuon alpinus): Asiatic wild dog",
    "- Habitat: Forests/grasslands of Central, East, SE Asia",
    "- Packs: 5-40, complex social structure, vocal communication",
    "- Cooperative hunting, hierarchical leadership",
    "",
    "How DIO maps dhole behavior to algorithm:",
    "- Lead vocalizer → Best solution guides search (leadership/exploitation)",
    "- Vocal signals → Vocalization influence (V) decays (explore→exploit)",
    "- Cooperative hunting → Pack adjusts together (avoid local optima)",
    "- Territorial instincts → Boundary constraints (feasible space)",
    "- Pack hierarchy → Multi-stage leadership transitions (adaptation)"
])

# Add sub-bullets
sub_points = [
    "Chase Alpha: Follow best hunter (exploitation)",
    "Random Pursuit: Chase random pack member (exploration)",
    "Pack Center: Move toward group average (cooperation)"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "Reference: El Romeh, Mirjalili, Šnel, Cluster Computing (2025), DOI: 10.1007/s10586-025-05543-2"
p.level = 0
p.font.size = Pt(18)

# ============================================================================
# SLIDE 5: DIO MATHEMATICAL FORMULATION
# ============================================================================
slide = add_content_slide(prs, "DIO Mathematical Formulation")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Position update combines three strategies:",
    "",
    "X_chase = X_alpha + r₁ × (X_alpha - X_i)",
    "X_random = X_r + r₂ × (X_r - X_i)",
    "X_scavenge = X_mean + r₃ × (X_mean - X_i)",
    "",
    "X_new = (X_chase + X_random + X_scavenge) / 3",
    "",
    "Where: r₁, r₂, r₃ ∈ [0,1] are random numbers"
])

# ============================================================================
# SLIDE 5A: DIO Key Equations & Biological Meaning
# ============================================================================
slide = add_content_slide(prs, "DIO: Key Equations & Biological Meaning")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "V(t) = 2 - (2*t/MaxIter): Vocalization influence (explore→exploit)",
    "B = V × r²: Movement scaling (aggressive→refined)",
    "C = r × sin(r): Sinusoidal oscillation (stochasticity)",
    "D_lead = C × (LeadVocalizer_pos - X)^2 + X^2: Adaptive distance",
    "X_new = LeadVocalizer_pos - B × sqrt(D_lead): Position update",
    "Boundary: If out of bounds, random reposition (territorial)",
    "Lead update: If fitness(X) > fitness(Lead), Lead = X (dynamic leadership)"
])

# ============================================================================
# SLIDE 5B: DIO ALGORITHM FLOWCHART
# ============================================================================
add_image_slide(
    prs,
    "DIO Algorithm Flowchart",
    os.path.join(parent_dir, "schemas and snippets", "dio_flowchart.png"),
    "Complete algorithmic flow: initialization, fitness evaluation, three hunting strategies, convergence"
)

# ============================================================================
# SLIDE 5C: DIO ADVANTAGES
# ============================================================================
slide = add_content_slide(prs, "DIO: Key Advantages")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "✅ Superior convergence speed (esp. unimodal functions): 15-20% faster than PSO, GWO",
    "✅ Excellent exploration-exploitation balance: avoids premature convergence",
    "✅ Outstanding local minima avoidance (multimodal functions)",
    "✅ Robust across problem types: unimodal, multimodal, composite, high-dimensional",
    "✅ Low variance, high stability: consistent results across 30 runs",
    "✅ Scalability: linear complexity, suitable for high-dimensional problems",
    "✅ Real-world engineering success: best results on welded beam & pressure vessel design",
    "✅ Biologically meaningful: genuine dhole pack inspiration, not just parameter tweaks",
    "✅ Easy implementation: straightforward, open-source code (GitHub, MATLAB)",
    "✅ Multi-stage leadership: dynamic adaptation, better than static GWO/PSO"
])

# ============================================================================
# SLIDE 5D: DIO LIMITATIONS
#+ ============================================================================
# SLIDE 5E: DIO vs. PSO, GWO, DE (Comparison)
# ============================================================================
slide = add_content_slide(prs, "DIO vs. PSO, GWO, DE: Comparison")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "DIO vs. PSO:",
    "- Inspiration: Dhole pack vs. bird flock",
    "- Leadership: Dynamic (multi-stage) vs. static (global best)",
    "- Parameter control: Vocalization (adaptive) vs. inertia (fixed)",
    "- Convergence: DIO faster, better balance; PSO can stagnate",
    "- Local minima: DIO excellent, PSO poor",
    "",
    "DIO vs. GWO:",
    "- Inspiration: Dhole vs. wolf pack",
    "- Leadership: Dynamic vs. static hierarchy",
    "- Movement: Sinusoidal+distance vs. linear encircling",
    "- Exploration: Vocalization decay vs. linear alpha",
    "- Boundary: Random reposition vs. clipping",
    "- Results: DIO better on F9 (2.2e-11 vs 0.31)",
    "",
    "DIO vs. DE:",
    "- Inspiration: Biological vs. mathematical",
    "- Mechanism: Cooperative vs. mutation/crossover",
    "- Parameter tuning: Moderate vs. high",
    "- Speed: DIO very fast, DE medium",
    "- Stability: DIO good, DE sometimes better",
    "- Head-to-head: DIO wins ~60% of tests, DE ~35%"
])
# ============================================================================
# SLIDE 5F: DIO Real-World Applications
# ============================================================================
slide = add_content_slide(prs, "DIO: Real-World Applications")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Engineering: Welded beam & pressure vessel design (cost reduction, safety)",
    "Machine learning: Hyperparameter tuning, feature selection (e.g., breast cancer diagnosis)",
    "Neural network weight optimization",
    "Resource allocation in networks",
    "Facility location, portfolio optimization",
    "Drug discovery, power systems, supply chain management",
    "Potential for any high-dimensional, single-objective optimization"
])
# ============================================================================
# SLIDE 5G: DIO Algorithm Pseudocode
# ============================================================================
slide = add_content_slide(prs, "DIO Algorithm: Pseudocode")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "1. Initialize population D with random solutions",
    "2. Evaluate fitness of all solutions",
    "3. For each iteration:",
    "   - Update vocalization influence: V = 2 - (2*iter/MaxIter)",
    "   - Find lead vocalizer (best solution)",
    "   - For each agent:",
    "       * r = random(0,1); B = V*r^2; C = r*sin(r)",
    "       * D_lead = C*(lead-X)^2 + X^2",
    "       * X_new = lead - B*sqrt(D_lead)",
    "       * If out of bounds: random reposition",
    "       * If fitness(X_new) > fitness(X): X = X_new",
    "   - If no improvement in N iters: break",
    "4. Return best solution found"
])
# ============================================================================
slide = add_content_slide(prs, "DIO: Key Limitations")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "⚠️ No formal theoretical convergence proof (empirical only)",
    "⚠️ Parameter sensitivity: requires tuning for V, C; not self-adaptive",
    "⚠️ Sometimes outperformed by Differential Evolution (DE) on some functions",
    "⚠️ Higher computational overhead per iteration (distance, sinusoidal calculations)",
    "⚠️ Population-based: not for single-solution or memory-constrained problems",
    "⚠️ Limited to single-objective optimization (no multi-objective support yet)",
    "⚠️ Tested mainly on benchmarks + 2 engineering problems; needs more real-world validation",
    "⚠️ Variable convergence on some fixed-dimension multimodal functions"
])

# ============================================================================
# SLIDE 6: NESTED OPTIMIZATION FRAMEWORK
# ============================================================================
slide = add_content_slide(prs, "Novel Nested Optimization Framework")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Two-level hierarchical optimization:",
])

p = tf.add_paragraph()
p.text = "Outer Loop: Hyperparameter Tuning"
p.level = 1
p.font.size = Pt(18)
p.font.bold = True

sub_points = [
    "Population: 3 dholes, 5 iterations",
    "Parameters: n_estimators, max_depth, min_samples_split, min_samples_leaf"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 2
    p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "Inner Loop: Feature Selection"
p.level = 1
p.font.size = Pt(18)
p.font.bold = True

sub_points = [
    "Population: 5 dholes, 10 iterations",
    "Selects optimal feature subset for each hyperparameter set"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 2
    p.font.size = Pt(16)

# ============================================================================
# SLIDE 6B: NESTED STRUCTURE VISUALIZATION
# ============================================================================
add_image_slide(
    prs,
    "Nested Optimization Architecture",
    os.path.join(parent_dir, "schemas and snippets", "shema4 (1).png"),
    "Two-level hierarchical structure: Outer loop (hyperparameters) + Inner loop (feature selection)"
)

# ============================================================================
# SLIDE 7: FITNESS FUNCTION
# ============================================================================
slide = add_content_slide(prs, "Multi-Objective Fitness Function")

# Add formula box
left = Inches(1.5)
top = Inches(2.5)
width = Inches(7)
height = Inches(1.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
p = text_frame.paragraphs[0]
p.text = "F = 0.99 × (1 - Accuracy) + 0.01 × (Features/Total)"
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Add explanation
left = Inches(1)
top = Inches(4.2)
width = Inches(8)
height = Inches(2)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
add_bullet_points(text_frame, [
    "Goal: Minimize F (balance accuracy and complexity)",
    "99% weight on accuracy, 1% on feature reduction",
    "Encourages Pareto-optimal solutions",
    "Favors fewer features when accuracy is comparable"
])

# ============================================================================
# SLIDE 7B: FITNESS FUNCTION SCHEMA
# ============================================================================
add_image_slide(
    prs,
    "Fitness Function: Balancing Accuracy and Complexity",
    os.path.join(parent_dir, "schemas and snippets", "shema5 (1).PNG"),
    "Multi-objective optimization: 99% accuracy weight + 1% feature reduction penalty"
)

# ============================================================================
# SLIDE 8: DATASET & EXPERIMENTAL SETUP
# ============================================================================
slide = add_content_slide(prs, "Dataset & Experimental Setup")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Dataset: Breast Cancer Wisconsin (Diagnostic)",
])

sub_points = [
    "569 samples (357 benign, 212 malignant)",
    "30 features from digitized cell nuclei images",
    "Standard benchmark in medical ML"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "Validation Strategy:"
p.level = 0
p.font.size = Pt(18)
p.font.bold = True

sub_points = [
    "30 independent runs with different train/test splits",
    "70/30 stratified split (random states: 42-71)",
    "Ensures statistical robustness and reproducibility"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

# ============================================================================
# SLIDE 9: OPTIMIZED CONFIGURATION
# ============================================================================
slide = add_content_slide(prs, "DIO-Optimized Configuration")

# Create table
left = Inches(1.5)
top = Inches(2)
width = Inches(7)
height = Inches(4)

# Add text boxes for results
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

p = text_frame.paragraphs[0]
p.text = "Optimized Hyperparameters:"
p.font.size = Pt(24)
p.font.bold = True

params = opt_results['best_hyperparameters']
param_text = [
    f"• n_estimators: {params['n_estimators']}",
    f"• max_depth: {params['max_depth']}",
    f"• min_samples_split: {params['min_samples_split']}",
    f"• min_samples_leaf: {params['min_samples_leaf']}"
]

for param in param_text:
    p = text_frame.add_paragraph()
    p.text = param
    p.font.size = Pt(20)
    p.space_after = Pt(8)

p = text_frame.add_paragraph()
p.text = f"\nSelected Features: {opt_results['selected_features']['count']} out of 30 (73% reduction!)"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(46, 204, 113)

# ============================================================================
# SLIDE 9B: CODE SNIPPET - DIO OPTIMIZE FUNCTION
# ============================================================================
add_image_slide(
    prs,
    "Implementation: DIO Optimize Function",
    os.path.join(parent_dir, "schemas and snippets", "dio_optimise_snippet.png"),
    "Core optimization loop with fitness evaluation and position updates"
)

# ============================================================================
# SLIDE 9C: CODE SNIPPET - FEATURE SELECTION OBJECTIVE
# ============================================================================
add_image_slide(
    prs,
    "Feature Selection Objective Function (RF)",
    os.path.join(parent_dir, "schemas and snippets", "feature_selection_objective_func_rf.png"),
    "Inner loop: Evaluates each feature subset with cross-validation"
)

# ============================================================================
# SLIDE 9D: CODE SNIPPET - HYPERPARAMETER OBJECTIVE
# ============================================================================
add_image_slide(
    prs,
    "Hyperparameter Objective Function (RF)",
    os.path.join(parent_dir, "schemas and snippets", "hyperparameter_objective_func_rf.png"),
    "Outer loop: Optimizes hyperparameters while calling feature selection"
)

# ============================================================================
# SLIDE 9E: CODE SNIPPET - OUTER OPTIMIZATION
# ============================================================================
add_image_slide(
    prs,
    "Outer Optimization Loop & Results Retrieval",
    os.path.join(parent_dir, "schemas and snippets", "outer_optimization_and_retreiving_results.png"),
    "Main execution: Runs nested optimization and extracts best configuration"
)

# ============================================================================
# SLIDE 10: SELECTED FEATURES
# ============================================================================
slide = add_content_slide(prs, "8 Selected Features")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

feature_names = opt_results['selected_features']['names']
add_bullet_points(tf, feature_names)

p = tf.add_paragraph()
p.text = "\nBalanced selection: Mean, Error, and Worst statistics"
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = RGBColor(52, 152, 219)

# ============================================================================
# SLIDE 11: OPTIMIZATION OVERFITTING DISCOVERY
# ============================================================================
add_image_slide(
    prs,
    "Major Discovery: Optimization Overfitting",
    os.path.join(parent_dir, "schemas and snippets", "Shema2 (1).png"),
    "Single-split optimization achieves 100% on one partition but generalizes poorly (94.72%). CV-based fixes this (96.26%)."
)

# ============================================================================
# SLIDE 12: RESULTS - VISUALIZATION
# ============================================================================
add_image_slide(
    prs,
    "Model Performance Comparison",
    "statistical_comparison_visualization.png",
    "30 independent runs across 10 machine learning models"
)

# ============================================================================
# SLIDE 12: RESULTS - THREE APPROACHES COMPARISON
# ============================================================================
slide = add_content_slide(prs, "Three Optimization Approaches Validated")

left = Inches(0.5)
top = Inches(2)
width = Inches(9)
height = Inches(4.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

results_text = """
Approach              Accuracy        Features  Opt.Time  Rank
───────────────────────────────────────────────────────────────
DIO-XGBoost          96.34% ±1.23%   17/30     54 sec    #1 🏆
(Single-Split)       (43% reduction)

DIO-RF-CV            96.26% ±1.33%   6/30      7.9 hrs   #3
(CV-Based)           (80% reduction)

DIO-RF-Single        94.72% ±1.41%   8/30      1 min     #7
(Initial)            (73% reduction)
───────────────────────────────────────────────────────────────

Key Insight: XGBoost achieves BEST OVERALL performance!
"""

p = text_frame.paragraphs[0]
p.text = results_text
p.font.name = 'Courier New'
p.font.size = Pt(16)

# ============================================================================
# SLIDE 13: KEY FINDINGS - UPDATED WITH ALL APPROACHES
# ============================================================================
slide = add_content_slide(prs, "What We Discovered: Algorithm-Dependent Behavior")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "🏆 DIO-XGBoost: BEST OVERALL (96.34%, 17 features, Rank #1)",
    "✓ Lowest standard deviation (1.23%) = Most stable across data partitions",
    "✓ Remarkably fast: 54-second optimization (526× faster than CV-RF!)",
    "✓ Significantly better than defaults (p=0.0426)",
    "",
    "🥉 DIO-RF-CV: Maximum Interpretability (96.26%, 6 features, Rank #3)",
    "✓ 80% feature reduction—highest among all configurations",
    "✓ Significantly better than defaults (p=0.0084)",
    "✓ Successfully addresses optimization overfitting",
    "",
    "💡 Unexpected Finding: Optimization Overfitting is Algorithm-Specific",
    "• RF single-split suffered overfitting (100% train → 94.72% generalization)",
    "• XGBoost single-split achieved top performance (98.83% → 96.34%)",
    "• Gradient boosting's built-in regularization eliminates need for expensive CV!"
])

# ============================================================================
# SLIDE 14: CV-BASED OPTIMIZATION VISUALIZATION
# ============================================================================
add_image_slide(
    prs,
    "CV-Based Optimization: Fixing Overfitting",
    os.path.join(parent_dir, "cv_optimization", "statistical_comparison_visualization_cv.png"),
    "96.26% ± 1.33% with 6 features (80% reduction) - Rank #3"
)

# ============================================================================
# SLIDE 15: XGBOOST OPTIMIZATION VISUALIZATION
# ============================================================================
add_image_slide(
    prs,
    "XGBoost Optimization: Best Overall Performance",
    os.path.join(parent_dir, "xgboost_statistical_comparison_visualization.png"),
    "96.34% ± 1.23% with 17 features (43% reduction) - Rank #1 🏆"
)

# ============================================================================
# SLIDE 15B: XGBOOST HYPERPARAMETER SEARCH SPACE (CANCER)
# ============================================================================
add_image_slide(
    prs,
    "XGBoost Hyperparameter Search Space (Medical)",
    os.path.join(parent_dir, "schemas and snippets", "xgboost_hyperparameters_search_space_cancer.png"),
    "5-dimensional search space: n_estimators, max_depth, learning_rate, subsample, colsample_bytree"
)

# ============================================================================
# SLIDE 16: STATISTICAL SIGNIFICANCE - ALL APPROACHES
# ============================================================================
slide = add_content_slide(prs, "Statistical Validation: All Three Approaches")

left = Inches(0.8)
top = Inches(2)
width = Inches(8.4)
height = Inches(4.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

p = text_frame.paragraphs[0]
p.text = "Wilcoxon Signed-Rank Test Results:"
p.font.size = Pt(22)
p.font.bold = True

test_results = [
    "",
    "DIO-XGBoost-Optimized (BEST - Rank #1):",
    "  • vs. XGBoost defaults: p=0.0426 (*) ✓",
    "  • vs. XGBoost (All): p=0.5067 (ns) - equivalent with 43% fewer features!",
    "  • vs. SVM/KNN/NB: p<0.001 (***) ✓✓✓",
    "",
    "DIO-RF-CV-Optimized (Rank #3):",
    "  • vs. RF defaults (6 feat): p=0.0084 (**) ✓✓",
    "  • vs. RF (All): p=0.0553 (ns) - comparable with 80% fewer features!",
    "  • vs. SVM: p<0.001 (***) ✓✓✓",
    "",
    "DIO-RF-Single (Rank #7 - optimization overfitting issue):",
    "  • vs. RF defaults (8 feat): p=0.165 (ns) - generalization problem",
    "  • vs. SVM/KNN: p<0.001 (***) ✓✓✓"
]

for result in test_results:
    p = text_frame.add_paragraph()
    p.text = result
    p.font.size = Pt(15)
    p.space_after = Pt(4)

# ============================================================================
# SLIDE 17: PARETO FRONTIER - THREE MODELS
# ============================================================================
slide = add_content_slide(prs, "Pareto Frontier: Three Validated Solutions")

left = Inches(0.8)
top = Inches(2)
width = Inches(8.4)
height = Inches(4.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

p = text_frame.paragraphs[0]
p.text = "Three Pareto-Optimal Solutions for Different Priorities:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(46, 204, 113)

points = [
    "",
    "🏆 Maximum Accuracy: DIO-XGBoost (96.34%, 17 features)",
    "  • Best overall performance across ALL experiments",
    "  • Lowest variance (1.23%) = Most stable predictions",
    "  • Fast optimization (54 seconds)",
    "  • Choose when: Accuracy is paramount",
    "",
    "🎯 Maximum Interpretability: DIO-RF-CV (96.26%, 6 features)",
    "  • Highest feature reduction (80%)",
    "  • Near-maximum accuracy with minimal complexity",
    "  • Choose when: Clinical transparency critical",
    "",
    "⚡ Rapid Prototyping: DIO-RF-Single (94.72%, 8 features)",
    "  • 1-minute optimization time",
    "  • Demonstrates optimization overfitting issue",
    "  • Choose when: Quick baseline needed"
]

for point in points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(16)
    p.space_after = Pt(4)

# ============================================================================
# SLIDE 18: PRACTICAL IMPLICATIONS - UPDATED
# ============================================================================
slide = add_content_slide(prs, "Clinical Deployment: Choose Your Model")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "🏆 Choose DIO-XGBoost (96.34%, 17 features) if:",
])

sub_points = [
    "Maximum accuracy is critical (e.g., high-stakes screening)",
    "Computational resources available",
    "43% feature reduction still provides efficiency gains"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "🎯 Choose DIO-RF-CV (96.26%, 6 features) if:"
p.level = 0
p.font.size = Pt(18)
p.font.bold = True

sub_points = [
    "Clinical interpretability paramount (only 6 measurements)",
    "Resource-constrained environments (80% cost reduction)",
    "Near-maximum accuracy (only 0.08% less than XGBoost)"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "⚡ Choose DIO-RF-Single (94.72%, 8 features) if:"
p.level = 0
p.font.size = Pt(18)
p.font.bold = True

sub_points = [
    "Rapid prototyping needed (1-minute optimization)",
    "Lower accuracy acceptable for initial screening"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

# ============================================================================
# SLIDE 19: EXTENSION TO COMPUTER VISION - CIFAR-10
# ============================================================================
slide = add_content_slide(prs, "Pushing the Limits: Can DIO Scale to Computer Vision?")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "🎯 Question: Does DIO work beyond medical tabular data?",
    "",
    "Dataset: CIFAR-10 (60K images, 10 classes)",
    "  • Full dataset: 50K train, 10K test",
    "  • Optimization subset: 2K train, 500 test (stratified)",
    "",
    "Feature Extraction: ResNet50 (pre-trained ImageNet)",
    "  • 2048-D deep learning features (68× larger than medical data!)",
    "  • Google Colab GPU extraction (~15 min)",
    "",
    "The challenge: High-dimensional optimization (30-D → 2048-D)",
    "  • Tests DIO scalability limits",
    "  • Validates algorithm selection methodology across domains"
])

# ============================================================================
# SLIDE 20: CIFAR-10 MODEL SELECTION
# ============================================================================
slide = add_content_slide(prs, "CIFAR-10: Model Selection Results")

left = Inches(2)
top = Inches(2.2)
width = Inches(6)
height = Inches(4)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

results_text = """
Full Dataset Comparison (50K train, 10K test):
═══════════════════════════════════════════════

Algorithm          Accuracy    Features
────────────────────────────────────────────
XGBoost            85.0%       2048    ✅
Random Forest      ~72%        2048
Logistic Reg.      ~60%        2048
KNN                ~57%        2048
────────────────────────────────────────────

Selection: XGBoost (best baseline performance)

Next: DIO optimization on 2K subset for feasibility
"""

p = text_frame.paragraphs[0]
p.text = results_text
p.font.name = 'Courier New'
p.font.size = Pt(16)

# ============================================================================
# SLIDE 21: CIFAR-10 DIO OPTIMIZATION RESULTS
# ============================================================================
slide = add_content_slide(prs, "CIFAR-10: DIO Optimization Results")

left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

p = text_frame.paragraphs[0]
p.text = "Optimization Subset (2K train, 500 test):"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(44, 62, 80)

results = [
    "",
    "📊 Performance:",
    "  • Baseline XGBoost: 80.8% accuracy",
    "  • DIO Optimized: 83.6% accuracy",
    "  • Improvement: +2.8% absolute (+3.47% relative)",
    "",
    "🎯 Feature Reduction:",
    "  • Original: 2,048 features",
    "  • Selected: 853 features",
    "  • Reduction: 58.35% (1,195 features eliminated!)",
    "  • Inference speedup: ~2.4×",
    "",
    "⚙️ Optimized Hyperparameters:",
    "  • n_estimators: 76 (was 100)",
    "  • max_depth: 5 (was 6)",
    "  • learning_rate: 0.217 (was 0.3)",
    "",
    "⏱️ Optimization Time: 5.4 hours (325 min)"
]

for result in results:
    p = text_frame.add_paragraph()
    p.text = result
    p.font.size = Pt(17)
    p.space_after = Pt(4)

# ============================================================================
# SLIDE 22: CROSS-DOMAIN COMPARISON
# ============================================================================
slide = add_content_slide(prs, "Cross-Domain Validation: Medical vs. Vision")

left = Inches(0.8)
top = Inches(2)
width = Inches(8.4)
height = Inches(4.8)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

comparison_text = """
Characteristic      Breast Cancer    CIFAR-10        Insight
═══════════════════════════════════════════════════════════════════
Domain              Medical          Computer Vision Multi-domain
Data Type           Tabular          Deep Features   Versatile
Features            30               2,048           68× larger
Classes             2 (binary)       10 (multi)      More complex
Training Samples    455              2,000           4.4× larger

Feature Reduction   80% (6 feat)     58.35% (853)    Both substantial
Accuracy Gain       +1.54% (CV)      +2.8%           Consistent
Optimization Time   7.9 hours        5.4 hours       Scalable
Algorithm           RF-CV            XGBoost         Task-specific

✅ Key Insight: DIO achieves substantial improvements in BOTH domains
   despite 68× dimensionality increase!
"""

p = text_frame.paragraphs[0]
p.text = comparison_text
p.font.name = 'Courier New'
p.font.size = Pt(14)

# ============================================================================
# SLIDE 22B: CROSS-DOMAIN FRAMEWORK SCHEMA
# ============================================================================
add_image_slide(
    prs,
    "Cross-Domain Validation Framework",
    os.path.join(parent_dir, "schemas and snippets", "shema1 (1).png"),
    "DIO framework validated across medical (30-D) and vision (2048-D) domains—68× dimensionality increase"
)

# ============================================================================
# SLIDE 22C: CROSS-DOMAIN RESULTS SCHEMA
# ============================================================================
add_image_slide(
    prs,
    "Cross-Domain Performance Summary",
    os.path.join(parent_dir, "schemas and snippets", "shema3 (1).png"),
    "Consistent improvement patterns: 58-80% feature reduction while maintaining/improving accuracy"
)

# ============================================================================
# SLIDE 23: CIFAR-10 PRACTICAL IMPLICATIONS
# ============================================================================
slide = add_content_slide(prs, "Real-World Impact: Why This Matters for Vision AI")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "🚀 Transfer Learning Made Better:",
    "  • Optimize frozen deep features without expensive CNN retraining",
    "  • 58.35% feature reduction while gaining +2.8% accuracy",
    "",
    "⚡ Edge Deployment Enabled:",
    "  • 2.4× faster inference (2048 → 853 features)",
    "  • Makes AI practical for smartphones, IoT devices, embedded systems",
    "",
    "💰 Cost-Effective for Startups:",
    "  • Only 4% of data needed (2K/50K samples)",
    "  • Valuable when annotation budgets are limited",
    "",
    "🔍 Surprising Discovery:",
    "  • 58% of ResNet50 features are redundant!",
    "  • Suggests efficient architecture design opportunities",
    "",
    "✅ Framework Transferability Proven:",
    "  • Medical (30-D) → Vision (2048-D) seamlessly",
    "  • Binary → Multi-class classification works"
])

# ============================================================================
# SLIDE 19: ALGORITHM VALIDATION
# ============================================================================
slide = add_content_slide(prs, "Framework Generalizability Validated")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "DIO Framework Successfully Applied to Multiple Algorithms:",
    "",
    "✓ Random Forest (Bagging ensemble)",
])

sub_points = [
    "Single-split: 94.72% (suffered optimization overfitting)",
    "CV-based: 96.26% (fixed overfitting, rank #3)"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "✓ XGBoost (Gradient boosting)"
p.level = 0
p.font.size = Pt(18)
p.font.bold = True

sub_points = [
    "Single-split: 96.34% (best overall, rank #1)",
    "Inherent regularization prevents optimization overfitting!"
]
for point in sub_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = ""
p.level = 0

p = tf.add_paragraph()
p.text = "Key Discovery: Algorithm-dependent optimization behavior"
p.level = 0
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(231, 76, 60)

# ============================================================================
# SLIDE 19B: EVOLUTION OF OPTIMIZATION APPROACHES
# ============================================================================
add_image_slide(
    prs,
    "Evolution of Optimization Strategies",
    os.path.join(parent_dir, "schemas and snippets", "Shema6 (1).PNG"),
    "Three approaches tested: Single-split RF → CV-based RF → Single-split XGBoost (best overall)"
)

# ============================================================================
# SLIDE 20: LIMITATIONS - UPDATED
# ============================================================================
slide = add_content_slide(prs, "Limitations & Lessons Learned")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "✓ RESOLVED: Optimization overfitting in Random Forest",
    "  → CV-based approach increased accuracy from 94.72% to 96.26%",
    "  → Feature reduction improved from 73% to 80%",
    "",
    "✓ DISCOVERY: Algorithm-dependent optimization behavior",
    "  → XGBoost's regularization enables successful single-split optimization",
    "  → Different algorithms require different validation strategies",
    "",
    "Remaining Limitations:",
    "• Single dataset evaluation (Breast Cancer Wisconsin only)",
    "• Limited hyperparameter spaces (4 RF, 5 XGBoost parameters)",
    "• No comparison with other metaheuristics (PSO, GA, ACO)",
    "• Feature selection stability not assessed across multiple runs"
])

# ============================================================================
# SLIDE 21: FUTURE WORK - UPDATED
# ============================================================================
slide = add_content_slide(prs, "Future Research Directions")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "1. Multi-dataset validation",
    "  → Lung cancer, diabetes, heart disease datasets",
    "  → Verify algorithm-dependent optimization findings",
    "",
    "2. CV-based XGBoost optimization",
    "  → Can 96.34% accuracy be improved further with CV?",
    "  → Compare single-split vs CV for gradient boosting",
    "",
    "3. Benchmark against other metaheuristics (PSO, GA, ACO, GWO)",
    "4. Extend to deep learning (neural architecture search)",
    "5. Feature selection stability analysis across multiple DIO runs",
    "6. Real-world clinical deployment and prospective validation",
    "7. Computational efficiency analysis and parallelization"
])

# ============================================================================
# SLIDE 22: CONTRIBUTIONS - UPDATED
# ============================================================================
slide = add_content_slide(prs, "Key Contributions")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "✓ First Python implementation of DIO algorithm",
    "✓ Novel nested optimization framework for simultaneous tuning",
    "✓ Three validated optimization approaches with trade-off analysis",
    "✓ Discovery of algorithm-dependent optimization overfitting",
    "✓ Multi-algorithm validation (Random Forest + XGBoost)",
    "✓ Rigorous statistical validation (30 independent runs per approach)",
    "✓ Comprehensive Pareto analysis with deployment recommendations",
    "✓ Open-source implementation for reproducible research",
    "✓ Best-in-class performance: 96.34% accuracy (Rank #1)"
])

# ============================================================================
# SLIDE 23: CONCLUSIONS - UPDATED
# ============================================================================
slide = add_content_slide(prs, "What We Learned")

left = Inches(0.8)
top = Inches(2)
width = Inches(8.4)
height = Inches(4.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

points = [
    "DIO framework successfully validated across TWO distinct domains:",
    "",
    "🏥 Medical Classification (Breast Cancer):",
    "  🏆 DIO-XGBoost: 96.34% accuracy, 17 features (Rank #1 BEST OVERALL)",
    "  🥉 DIO-RF-CV: 96.26% accuracy, 6 features (80% reduction)",
    "",
    "🖼️ Computer Vision (CIFAR-10 Images):",
    "  • Full dataset baseline: 85% (XGBoost, 2048 features)",
    "  • DIO optimized: 83.6% on subset (+2.8%, 58.35% reduction)",
    "  • 2.4× inference speedup for edge deployment",
    "",
    "🔑 Unexpected Discoveries:",
    "  • Optimization overfitting is algorithm-dependent (not universal!)",
    "  • DIO scales remarkably: 30-D → 2048-D (68× increase)",
    "  • Even deep learning features have massive redundancy (58%)",
    "",
    "✅ Multi-domain framework proven for medical AI & computer vision"
]

for i, point in enumerate(points):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.space_after = Pt(6)


# ============================================================================
# SLIDE 24: THANK YOU / Q&A
# ============================================================================
slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# Title
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(1)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
p = text_frame.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(60)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(44, 62, 80)

# Subtitle with achievement
left = Inches(1)
top = Inches(3.2)
width = Inches(8)
height = Inches(1.2)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
p = text_frame.paragraphs[0]
p.text = "🏆 Best Medical Performance: 96.34% Accuracy"
p.font.size = Pt(26)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(46, 204, 113)
p.font.bold = True

p = text_frame.add_paragraph()
p.text = "✅ Cross-Domain Validation: Medical + Vision"
p.font.size = Pt(24)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(52, 152, 219)
p.space_before = Pt(8)

p = text_frame.add_paragraph()
p.text = "Questions & Discussion"
p.font.size = Pt(32)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(44, 62, 80)
p.space_before = Pt(12)

# Contact/Links
left = Inches(1)
top = Inches(5.8)
width = Inches(8)
height = Inches(1)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

contact_info = [
    "GitHub: amine-dubs/dio-optimization",
    "Domains: Medical Classification (96.34%) + Computer Vision (83.6%)",
    "Datasets: UCI Breast Cancer + CIFAR-10 ResNet50 Features"
]

for i, info in enumerate(contact_info):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = info
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 25: REFERENCES
# ============================================================================
slide = add_content_slide(prs, "References")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "[1] El Romeh, A., Mirjalili, S., Šnel, V. (2025). Dholes-Inspired Optimization (DIO). Cluster Computing, Springer. DOI: 10.1007/s10586-025-05543-2. Open-source: github.com/AlyromehDholes-Inspired-Optimization-DIO, MathWorks File Exchange. https://link.springer.com/article/10.1007/s10586-025-05543-2",
    "[2] UCI Machine Learning Repository: Breast Cancer Wisconsin (Diagnostic) Data Set.",
    "[3] Krizhevsky, A. (2009). CIFAR-10 Dataset.",
    "[4] Chen, T., & Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System.",
    "[5] Breiman, L. (2001). Random Forests. Machine Learning.",
    "[6] scikit-learn: Machine Learning in Python.",
    "[7] python-pptx documentation."
])

# ============================================================================
# SAVE PRESENTATION
# ============================================================================
output_file = os.path.join(script_dir, "DIO_Research_Presentation.pptx")
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"⏱️  Estimated presentation time: ~20-22 minutes")
print(f"\n🎯 Highlights:")
print(f"   🏆 DIO-XGBoost (Medical): 96.34% accuracy (Rank #1)")
print(f"   🥉 DIO-RF-CV (Medical): 96.26% accuracy, 6 features (80% reduction)")
print(f"   �️ DIO-XGBoost (CIFAR-10): 83.6% accuracy, 58.35% feature reduction")
print(f"   💡 Cross-domain validation: 30-D → 2048-D (68× scale-up)")
print(f"\n📝 Next steps:")
print(f"   1. Open {output_file} in PowerPoint")
print(f"   2. Review medical + vision results")
print(f"   3. Add speaker notes if needed")
print(f"   4. Practice timing (~45-50 seconds per slide)")
print(f"   5. Ready to present multi-domain research!")
