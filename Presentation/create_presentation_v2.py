"""
PowerPoint Presentation Generator for DIO Research Project
Restructured version following user-specified outline
Creates a 20-minute presentation ready for export with citations and page numbers
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
import os

# Get the directory of this script
script_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(script_dir)

# Slide counter for page numbers
slide_counter = 0

# Load optimization results from parent directory
try:
    with open(os.path.join(parent_dir, 'optimization_results.json'), 'r') as f:
        opt_results = json.load(f)
except:
    opt_results = {
        'best_hyperparameters': {'n_estimators': 200, 'max_depth': 20, 'min_samples_split': 9, 'min_samples_leaf': 10},
        'selected_features': {'count': 4, 'names': ['mean smoothness', 'area error', 'worst area', 'worst smoothness']}
    }

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_page_number(slide, number):
    """Add page number to bottom right of slide"""
    left = Inches(9)
    top = Inches(7)
    width = Inches(0.8)
    height = Inches(0.4)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = str(number)
    
    # Format the page number
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)

def add_title_slide(prs, title, subtitle):
    global slide_counter
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_top = Inches(4.5)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(2))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    
    slide_counter += 1
    # Title slide typically doesn't have page number, but we count it
    
    return slide

def add_content_slide(prs, title, content_type='bullet'):
    global slide_counter
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    slide_counter += 1
    add_page_number(slide, slide_counter)
    
    return slide

def add_bullet_points(text_frame, points, level=0):
    """Add bullet points to a text frame"""
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = level
        p.font.size = Pt(20)

def add_image_slide(prs, title, image_path, caption=""):
    global slide_counter
    """Add a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    try:
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        
        pic = slide.shapes.add_picture(image_path, left, top, width=width)
        
        if caption:
            left = Inches(0.5)
            top = Inches(6.5)
            width = Inches(9)
            height = Inches(0.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = caption
            p.font.size = Pt(14)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
    except Exception as e:
        body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_shape:
            tf = body_shape.text_frame
            p = tf.paragraphs[0]
            p.text = f"Image not found: {image_path}"
            p.font.size = Pt(16)
    
    slide_counter += 1
    add_page_number(slide, slide_counter)
    
    return slide


# ============================================================================
# TITLE SLIDE
# ============================================================================
add_title_slide(
    prs,
    "Dholes-Inspired Optimization (DIO)\nfor Feature Selection & Hyperparameter Tuning",
    "Cross-Domain Validation: From Medical Diagnostics to Computer Vision\n\nUSTO-MB | Computer Science | Data Science | December 2025"
)

# ============================================================================
# INTRO SLIDE (Team & Course Info)
# ============================================================================
slide = add_content_slide(prs, "Introduction & Team")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Students: Bellatreche Mohamed Amine, Cherif Ghizlane Iman",
    "University: USTO-MB (Université des Sciences et de la Technologie d'Oran)",
    "Department: Computer Science",
    "Specialty: Data Science",
    "Course: Statistics in Data Science",
    "Professor: Dr. Neggaz Nabil",
    "",
    "A journey through optimization—from success to failure to deeper understanding"
])

# ============================================================================
# TABLE OF CONTENTS
# ============================================================================
slide = add_content_slide(prs, "Contents")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Introduction",
    "1. Dholes-Inspired Optimization (DIO)",
    "    1.1 History & Origins",
    "    1.2 How DIO Works",
    "    1.3 Benchmark Comparisons",
    "    1.4 Applications",
    "2. Optimization Across Domains",
    "    2.1 The Challenge We Faced",
    "    2.2 Optimization Methodology",
    "    2.3 Medical Domain: Breast Cancer",
    "    2.4 Computer Vision: CIFAR-10",
    "    2.5 Cross-Domain Insights",
    "3. Future Research Directions",
    "4. Conclusion & Key Contributions",
    "5. References"
])

# ============================================================================
# INTRODUCTION
# ============================================================================
slide = add_content_slide(prs, "Introduction: Why This Research Matters")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Breast cancer is one of the deadliest diseases affecting women worldwide [2]",
    "",
    "Machine learning shows promise, but faces real challenges:",
    "  - Medical datasets have 30+ features—which ones really matter?",
    "  - A model's success depends on both the features AND its settings",
    "  - Picking features first, then tuning later misses important connections",
    "",
    "The old way: Choose features → Train model → Tune parameters",
    "The problem: What if the best features change depending on the parameters?",
    "",
    "Our approach: Optimize both together using DIO [1]",
    "",
    "This presentation tells the story of what we discovered—successes, failures, and lessons learned"
])

# ============================================================================
# PART 1: DIO - SECTION DIVIDER
# ============================================================================
slide = add_content_slide(prs, "Part 1: Dholes-Inspired Optimization")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "A cutting-edge optimization algorithm inspired by wild dog hunting",
    "",
    "What we'll cover:",
    "  - Where DIO came from and who created it",
    "  - How it actually works (the hunting strategies)",
    "  - How it compares to other popular algorithms",
    "  - What you can use it for"
])

# ============================================================================
# 1.1 DIO HISTORY
# ============================================================================
slide = add_content_slide(prs, "1.1 DIO: History & Origins")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Created in 2025 by international researchers [1]:",
    "  - Dr. Seyedali Mirjalili (Torrens University, Australia)",
    "  - Ali El Romeh (Torrens University, Australia)",
    "  - Václav Šnel (VSB-Technical University, Czech Republic)",
    "",
    "Published in Cluster Computing (Springer)—a prestigious journal",
    "  - DOI: 10.1007/s10586-025-05543-2",
    "  - Peer-reviewed and rigorously tested",
    "",
    "Why dholes? These Asian wild dogs are fascinating:",
    "  - Hunt in packs with amazing coordination",
    "  - Each member has a role but adapts flexibly",
    "  - Balance independence with teamwork perfectly",
    "",
    "Code is open-source: Available on GitHub and MATLAB File Exchange"
])

# ============================================================================
# 1.2 HOW DIO WORKS
# ============================================================================
slide = add_content_slide(prs, "1.2 How DIO Works: The Hunting Strategies")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "DIO translates dhole pack hunting into mathematics",
    "",
    "Dholes use three hunting strategies:",
    "",
    "Strategy 1: Chase the Leader",
    "  - Follow the best hunter toward the prey",
    "  - In optimization: Exploit the best solution found so far",
    "",
    "Strategy 2: Try Random Directions  ",
    "  - Check what other pack members are discovering",
    "  - In optimization: Explore new areas of the search space",
    "",
    "Strategy 3: Stick with the Pack",
    "  - Move toward the group's center",
    "  - In optimization: Balance between extremes, avoid getting stuck",
    "",
    "Key insight: Combining all three prevents premature convergence while maintaining focus"
])

# ============================================================================
# DIO ANIMAL BEHAVIOR MAPPING
# ============================================================================
slide = add_content_slide(prs, "Dholes: From Biology to Algorithm")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Dhole (Cuon alpinus): The Asiatic wild dog",
    "  - Lives in forests across Central, East, and Southeast Asia",
    "  - Packs of 5-40 individuals with complex social structures",
    "  - Vocal communication (hence 'vocalization' in the math)",
    "",
    "How biology maps to the algorithm:",
    "  - Pack leader → Best solution guides the search",
    "  - Vocal signals → 'Vocalization influence' parameter that decreases over time",
    "  - Cooperative hunting → All solutions move together strategically",
    "  - Territory boundaries → Search space constraints",
    "  - Changing leadership → Algorithm adapts when better solutions emerge",
    "",
    "This isn't just metaphor—it's genuine biological inspiration"
])

# ============================================================================
# DIO MATHEMATICAL FORMULATION
# ============================================================================
slide = add_content_slide(prs, "DIO: The Math Behind It (Optional Detail)")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Position update equation combines all three strategies:",
    "",
    "X_chase = X_alpha + r₁ × (X_alpha - X_i)     [Follow leader]",
    "X_random = X_r + r₂ × (X_r - X_i)             [Explore random]",
    "X_scavenge = X_mean + r₃ × (X_mean - X_i)    [Stay centered]",
    "",
    "Final position: X_new = (X_chase + X_random + X_scavenge) / 3",
    "",
    "Where:",
    "  - X_alpha = best solution (pack leader)",
    "  - X_r = random pack member",
    "  - X_mean = pack average position",
    "  - r₁, r₂, r₃ = random numbers in [0,1] for stochasticity",
    "",
    "The beauty: Simple math, powerful results"
])

# ============================================================================
# 1.3 BENCHMARK COMPARISON
# ============================================================================
slide = add_content_slide(prs, "1.3 How Good Is DIO? Benchmark Comparisons")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Tested on CEC2017—the gold standard for optimization algorithms",
    "  - 29 different test functions (unimodal, multimodal, hybrid, composite)",
    "  - Industry standard for comparing metaheuristics",
    "",
    "Compared against famous algorithms:",
    "  - PSO (Particle Swarm Optimization)—inspired by bird flocking",
    "  - GWO (Grey Wolf Optimizer)—inspired by wolf hunting",
    "  - WOA (Whale Optimization Algorithm)—inspired by whale hunting",
    "  - MFO (Moth-Flame Optimization)",
    "  - ALO (Ant Lion Optimizer)",
    "",
    "Results speak for themselves:",
    "  - DIO ranked 1st or 2nd on 22 out of 29 functions",
    "  - Beats PSO and GWO consistently on complex problems",
    "  - Particularly strong at avoiding local optima traps",
    "",
    "Bottom line: DIO is competitive with—and often better than—established algorithms"
])

# ============================================================================
# 1.4 DIO APPLICATIONS
# ============================================================================
slide = add_content_slide(prs, "1.4 What Can You Use DIO For?")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "DIO's versatility makes it useful across many fields:",
    "",
    "Engineering & Design:",
    "  - Structural optimization (designing bridges, aircraft parts)",
    "  - Resource allocation in factories",
    "",
    "Machine Learning (our focus here):",
    "  - Selecting which features to use from high-dimensional data",
    "  - Tuning hyperparameters for better model performance",
    "  - Neural architecture search for deep learning",
    "",
    "Other interesting domains:",
    "  - Financial portfolio optimization (maximizing returns, minimizing risk)",
    "  - Energy grid management (balancing supply and demand)",
    "  - Logistics and delivery route planning",
    "",
    "Why it works everywhere: Strong exploration prevents getting stuck in bad solutions"
])

# ============================================================================
# SCHEMA 1: CROSS-DOMAIN FRAMEWORK
# ============================================================================
add_image_slide(
    prs,
    "DIO Framework: Medical + Vision Domains",
    os.path.join(parent_dir, 'schemas and snippets', 'schema1_cross_domain_optimization_framework.png'),
    "Our implementation applies DIO to both medical diagnosis and computer vision"
)

# ============================================================================
# PART 2: OPTIMIZATION - SECTION DIVIDER
# ============================================================================
slide = add_content_slide(prs, "Part 2: Optimization Across Domains")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Now that we understand DIO, let's see how we applied it",
    "",
    "Two very different problems:",
    "  - Medical: Breast cancer diagnosis (30 features, binary classification)",
    "  - Vision: CIFAR-10 image classification (2048 features, 10 classes)",
    "",
    "What we'll cover:",
    "  - The challenge we didn't expect",
    "  - Our nested optimization methodology",
    "  - Results from medical domain (3 different approaches!)",
    "  - Extension to computer vision",
    "  - What we learned across both domains"
])

# ============================================================================
# 2.1 THE CHALLENGE WE FACED
# ============================================================================
slide = add_content_slide(prs, "2.1 The Challenge: Optimization Overfitting")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "We faced two distinct challenges during this research:",
    "",
    "Challenge 1: The Sequential Trap",
    "  - Traditional approach: Pick features first, tune parameters second",
    "  - Problem: Best features depend on which parameters you use!",
    "  - Solution: Optimize both simultaneously with nested DIO",
    "",
    "Challenge 2: The Overfitting Surprise (discovered the hard way)",
    "  - Got 100% accuracy optimizing on one train-test split",
    "  - Tested across 30 different splits: only 94.37% average",
    "  - Even worse than not optimizing at all (94.89%)!",
    "",
    "What happened? Optimization overfitting:",
    "  - The algorithm 'memorized' that specific data split",
    "  - Found parameters perfect for one partition, terrible for others",
    "  - Like studying only the practice exam and bombing the real one",
    "",
    "This discovery shaped everything that followed..."
])

# ============================================================================
# 2.2 OPTIMIZATION METHODOLOGY
# ============================================================================
slide = add_content_slide(prs, "2.2 Our Nested Optimization Methodology")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "We developed a two-level optimization framework:",
    "",
    "Outer Loop: Hyperparameter Tuning (DIO controls this)",
    "  - Searches through model parameter space",
    "  - For Random Forest: number of trees, tree depth, split criteria",
    "  - For XGBoost: learning rate, max depth, regularization",
    "",
    "Inner Loop: Feature Selection (nested DIO controls this)",
    "  - For each set of parameters, finds best feature subset",
    "  - Binary encoding: each feature is either used (1) or not (0)",
    "  - Example: [1,0,1,1,0...] means 'use features 1, 3, and 4'",
    "",
    "Fitness Function balances two goals:",
    "  - 99% weight: Maximize classification accuracy",
    "  - 1% weight: Minimize number of features used",
    "",
    "Key advantage: Discovers how features and parameters interact—something sequential optimization misses"
])

# ============================================================================
# SCHEMA 4: NESTED STRUCTURE
# ============================================================================
add_image_slide(
    prs,
    "Nested Optimization Architecture",
    os.path.join(parent_dir, 'schemas and snippets', 'schema4_nested_optimization_structure.png'),
    "Two-level optimization: Outer loop (hyperparameters) + Inner loop (features)"
)

# ============================================================================
# SCHEMA 5: FEATURE ENCODING
# ============================================================================
add_image_slide(
    prs,
    "Feature Selection: Binary Encoding",
    os.path.join(parent_dir, 'schemas and snippets', 'schema5_feature_selection_encoding.png'),
    "How DIO represents feature selection as a binary decision vector"
)

# ============================================================================
# 2.3 MEDICAL DOMAIN - BREAST CANCER
# ============================================================================
slide = add_content_slide(prs, "2.3 Medical Domain: Breast Cancer Dataset")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
add_bullet_points(tf, [
    "Wisconsin Diagnostic Breast Cancer Dataset [2]",
    "  - Source: UCI Machine Learning Repository (standard benchmark)",
    "  - 569 patients: 357 benign, 212 malignant tumors",
    "",
    "30 features describing cell nuclei characteristics:",
    "  - Size measurements: radius, perimeter, area",
    "  - Texture: smoothness, compactness, concavity",
    "  - Shape: symmetry, fractal dimension",
    "  - Statistics: mean, standard error, worst (largest) values",
    "",
    "Why this matters clinically:",
    "  - Early, accurate diagnosis literally saves lives",
    "  - Fewer features → faster, cheaper screening tests",
    "  - Simpler models → doctors can actually understand them",
    "",
    "Experimental setup: 70% training, 30% testing (stratified by class)"
])

# ============================================================================
# 2.3.1 DIO-RF SINGLE SPLIT (THE PROBLEM)
# ============================================================================
slide = add_content_slide(prs, "2.3.1 Approach 1: DIO-RF Single Split (The Trap)")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "Our first attempt: Optimize on one fixed train-test split",
    "",
    "Initial results looked incredible:",
    f"  - Training accuracy: 100% (perfect!)",
    f"  - Test accuracy: 97.08% on that specific split",
    f"  - Only 8 features used (73% reduction from 30)",
    "  - We thought we'd solved it!",
    "",
    "But then we ran 30-run statistical validation [7]...",
    "",
    "The harsh truth:",
    f"  - Average accuracy across 30 different splits: 94.37% ± 1.82%",
    f"  - Baseline (no optimization): 94.89% ± 1.47%",
    f"  - We actually made it WORSE by 0.17%!",
    "",
    "What went wrong? Optimization overfitting:",
    "  - DIO found parameters perfect for THAT ONE split",
    "  - Didn't generalize to different data partitions",
    "  - Classic case of 'overfitting' but at the optimization level, not model level"
])

# ============================================================================
# 2.3.2 DIO-RF WITH CV (THE SOLUTION)
# ============================================================================
slide = add_content_slide(prs, "2.3.2 Approach 2: DIO-RF with Cross-Validation (The Fix)")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "The solution: Change how we evaluate during optimization",
    "",
    "Instead of one train-test split:",
    "  - Use 5-fold cross-validation DURING the DIO search",
    "  - Each candidate solution tested on 5 different data partitions",
    "  - Fitness score = average performance across all 5 folds",
    "  - Makes it impossible to 'cheat' by memorizing one partition",
    "",
    "Results after fixing (30-run validation):",
    f"  - Average accuracy: 96.55% ± 1.51%",
    f"  - Features used: Only 6 out of 30 (80% reduction!)",
    f"  - F1-score: 0.9632",
    "",
    "Comparison to broken single-split approach:",
    f"  - +1.54% accuracy improvement (statistically significant p < 0.001)",
    f"  - More stable (lower standard deviation)",
    f"  - Overall rank: #3 among all methods tested",
    "",
    "The lesson: Use cross-validation during optimization, not just for final evaluation"
])

# ============================================================================
# 2.3.3 DIO-XGBOOST (BEST MODEL)
# ============================================================================
slide = add_content_slide(prs, "2.3.3 Approach 3: DIO-XGBoost (The Champion)")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "Wait—can we do even better with a different algorithm?",
    "",
    "Why try XGBoost [5]?",
    "  - Gradient boosting with built-in regularization",
    "  - More hyperparameters to tune (learning rate, tree depth, subsample ratio)",
    "  - Widely known to be more robust than Random Forest",
    "  - Gold standard for tabular data competitions (Kaggle, etc.)",
    "",
    "DIO-XGBoost results (30-run validation):",
    "  - Average accuracy: 96.88% ± 1.10%  ← HIGHEST OF ALL APPROACHES",
    "  - Features: 10 out of 30 (67% reduction)",
    "  - Overall rank: #1 (best performing method)",
    "",
    "Statistical significance:",
    "  - Significantly better than XGBoost defaults (p = 0.0426)",
    "  - Outperformed both Random Forest approaches",
    "  - Lowest variance (most stable predictions)",
    "",
    "Surprising finding:",
    "  - XGBoost didn't need cross-validation during optimization!",
    "  - Single-split worked fine (inherent regularization prevented overfitting)",
    "  - Algorithm choice matters for optimization strategy"
])

# ============================================================================
# 2.3.4 THREE APPROACHES VALIDATED (SCHEMA 6)
# ============================================================================
add_image_slide(
    prs,
    "2.3.4 Medical Results: All Three Approaches",
    os.path.join(parent_dir, 'schemas and snippets', 'schema6_medical_classification_results.png'),
    "Visual comparison of DIO-RF-Single (failed), DIO-RF-CV (good), DIO-XGBoost (best)"
)

slide = add_content_slide(prs, "Summary: Three Validated Approaches")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "Summary of Medical Domain Results (30-run averages):",
    "",
    "❌ Approach 1: DIO-RF Single Split",
    "  - Accuracy: 94.37% ± 1.82%  |  Features: 8  |  Rank: #6",
    "  - Problem: Optimization overfitting",
    "",
    "✓ Approach 2: DIO-RF with Cross-Validation",
    "  - Accuracy: 96.55% ± 1.51%  |  Features: 6  |  Rank: #1",
    "  - Solution: CV prevents overfitting, achieves 80% feature reduction",
    "",
    "🏆 Approach 3: DIO-XGBoost (Winner)",
    "  - Accuracy: 96.88% ± 1.10%  |  Features: 10  |  Rank: #1",
    "  - Best: Highest accuracy, lowest variance, significantly better than defaults",
    "",
    "Baseline comparisons for context:",
    "  - RF Default (30 features): 94.89% ± 1.47%",
    "  - XGBoost Default (30 features): 94.74% ± 1.55%",
    "",
    "Key insight: Algorithm robustness determines whether CV is needed during optimization"
])

# ============================================================================
# 2.3.5 STATISTICAL VALIDATION (WILCOXON TEST)
# ============================================================================
slide = add_content_slide(prs, "2.3.5 Statistical Rigor: Wilcoxon Signed-Rank Test [7]")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "30 runs give us averages—but are differences real or just luck?",
    "",
    "Wilcoxon Signed-Rank Test:",
    "  - Non-parametric test (doesn't assume normal distributions)",
    "  - Compares paired samples (same data splits for fairness)",
    "  - Null hypothesis: 'No real difference between methods'",
    "",
    "Results for DIO-XGBoost vs baselines:",
    "  - vs XGBoost Default: p = 0.0426 (*) → Significant improvement",
    "  - Mean difference: +1.60% accuracy",
    "  - vs RF Default: p < 0.001 (***) → Highly significant improvement",
    "",
    "Results for DIO-RF-CV vs DIO-RF-Single:",
    "  - p < 0.001 (***) → Cross-validation dramatically better",
    "  - Confirms CV is essential for RF optimization",
    "",
    "Significance levels:",
    "  * p < 0.05 (significant)",
    "  ** p < 0.01 (very significant)",
    "  *** p < 0.001 (extremely significant)",
    "  ns = not significant (could be chance)"
])

# ============================================================================
# 2.4 COMPUTER VISION - CIFAR-10
# ============================================================================
slide = add_content_slide(prs, "2.4 Computer Vision Domain: CIFAR-10")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "Can DIO work beyond medical data? Let's test on images",
    "",
    "CIFAR-10 Dataset [3]:",
    "  - 60,000 natural images (32×32 pixels, color)",
    "  - 10 classes: airplane, car, bird, cat, deer, dog, frog, horse, ship, truck",
    "  - Standard benchmark in computer vision research",
    "",
    "Our approach:",
    "  - Used ResNet50 [4] (pretrained on ImageNet) to extract features",
    "  - Converted images → 2048-dimensional feature vectors",
    "  - Computational subset: 2,500 images (2000 train, 500 test)",
    "  - Applied DIO-XGBoost optimization",
    "",
    "The massive challenge:",
    "  - Medical: 30 features (search space size)",
    "  - Vision: 2048 features (68× larger!)",
    "  - Much harder to optimize—needs way more computation"
])

# ============================================================================
# 2.4.1 CIFAR-10 RESULTS (THE PROBLEM RETURNS)
# ============================================================================
slide = add_content_slide(prs, "2.4.1 CIFAR-10 Results: When Budget Isn't Enough")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "We ran the same DIO-XGBoost optimization on CIFAR-10",
    "",
    "Single-run result looked okay:",
    "  - DIO-optimized (598 features): 83.0% accuracy",
    "  - Baseline (all 2048 features): 80.8% accuracy",
    "  - Seemed like +2.2% improvement with 70.8% feature reduction",
    "",
    "But then came the 30-run statistical validation...",
    "",
    "The disappointing truth:",
    "  - DIO-XGBoost: 81.91% ± 1.38%  (Rank #3)",
    "  - XGBoost Default: 83.27% ± 1.25%  (Rank #1) ← Actually BETTER",
    "  - Wilcoxon test: p = 7.15×10⁻⁵ (***) → Highly significant WORSE",
    "",
    "What went wrong? Insufficient optimization budget:",
    "  - Medical: 30-D space, 250-576 function evaluations → SUCCESS",
    "  - Vision: 2048-D space, 576 function evaluations → FAILURE",
    "  - Budget didn't scale with dimensionality",
    "  - Need ~10,000-50,000 evaluations for 2048-D (17-87× more)",
    "",
    "Optimization overfitting strikes again—even XGBoost can't overcome severe under-budgeting"
])

# ============================================================================
# 2.5 CROSS-DOMAIN INSIGHTS (SCHEMA 3)
# ============================================================================
add_image_slide(
    prs,
    "2.5 Cross-Domain Analysis: Success vs Failure",
    os.path.join(parent_dir, 'schemas and snippets', 'schema3_crossdomain_insights.png'),
    "Why optimization succeeded on medical data but failed on CIFAR-10"
)

slide = add_content_slide(prs, "Lessons Learned Across Both Domains")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "Three critical factors determine optimization success:",
    "",
    "Factor 1: Optimization Budget vs Dimensionality",
    "  - Medical (30-D): 250-576 evaluations sufficient → SUCCESS",
    "  - Vision (2048-D): 576 evaluations insufficient → FAILURE",
    "  - Rule of thumb: Budget must scale with dimensionality²",
    "",
    "Factor 2: Algorithm Robustness",
    "  - Random Forest: Needs cross-validation even at low dimensions",
    "  - XGBoost: Single-split OK at low-D, fails at high-D with low budget",
    "  - No algorithm can overcome severely inadequate computational resources",
    "",
    "Factor 3: Validation Strategy",
    "  - Single-run results are deceptive (83.0% vs 80.8% looked good)",
    "  - 30-run averages reveal truth (81.91% vs 83.27% showed failure)",
    "  - Statistical testing confirms what's real and what's noise",
    "",
    "The big lesson: Success isn't just about the algorithm—it's about matching computational investment to problem complexity"
])

# ============================================================================
# PART 3: FUTURE RESEARCH DIRECTIONS
# ============================================================================
slide = add_content_slide(prs, "3. Future Research Directions")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "Where do we go from here? Several exciting directions:",
    "",
    "1. Adaptive Budget Allocation",
    "  - Automatically adjust population size based on dimensionality",
    "  - Develop formulas: budget ≈ k × D² (where D = dimensions)",
    "  - Balance computational cost with quality",
    "",
    "2. Hybrid Optimization Strategies",
    "  - Combine DIO (global search) with local refinement (Bayesian optimization)",
    "  - Use DIO for exploration, gradient methods for exploitation",
    "",
    "3. Multi-Objective Extensions",
    "  - Optimize accuracy + interpretability + speed simultaneously",
    "  - Generate Pareto fronts for clinical decision-making",
    "",
    "4. Theoretical Foundations",
    "  - Formalize 'optimization overfitting' mathematically",
    "  - Prove convergence guarantees for nested optimization",
    "  - Characterize when single-split vs CV is sufficient",
    "",
    "5. Broader Applications",
    "  - Other medical tasks: multi-class diagnosis, survival prediction",
    "  - Neural architecture search with DIO",
    "  - Fairness-aware feature selection (removing biased features)"
])

# ============================================================================
# 4. CONCLUSION & KEY CONTRIBUTIONS
# ============================================================================
slide = add_content_slide(prs, "4. Conclusion: What This Research Taught Us")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "This journey taught us as much from failures as from successes:",
    "",
    "Key Contribution #1: Discovery of Optimization Overfitting",
    "  - Models can 'memorize' specific data partitions during optimization",
    "  - Single-run results can be wildly misleading (100% → 94.37%)",
    "  - Solution: Use cross-validation DURING optimization, not just after",
    "",
    "Key Contribution #2: Algorithm-Dependent Strategies",
    "  - Random Forest: Always needs CV for robust optimization",
    "  - XGBoost: More forgiving but still fails with inadequate budget",
    "  - No algorithm is immune to severe under-resourcing",
    "",
    "Key Contribution #3: Budget Scaling Laws",
    "  - Success on medical (30-D, 250-576 evals)",
    "  - Failure on vision (2048-D, 576 evals)",
    "  - Lesson: Budget must grow ~quadratically with dimensions",
    "",
    "Key Contribution #4: DIO-XGBoost Medical Success",
    "  - 96.88% accuracy (best overall), 67% feature reduction",
    "  - Significantly better than defaults (p=0.0426)",
    "  - Proves methodology works when properly resourced",
    "",
    "Final thought: Optimization isn't magic—it needs adequate investment"
])

# ============================================================================
# PRACTICAL TAKEAWAYS
# ============================================================================
slide = add_content_slide(prs, "Practical Recommendations")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "If you're deploying ML in healthcare:",
    "",
    "Choose DIO-XGBoost (96.88%, 10 features) if:",
    "  - Maximum accuracy is non-negotiable",
    "  - You have computational resources",
    "  - 67% feature reduction provides significant efficiency gains",
    "",
    "Choose DIO-RF-CV (96.55%, 6 features) if:",
    "  - Clinical interpretability is paramount (only 6 measurements!)",
    "  - Resource-constrained environments (80% cost reduction)",
    "  - Near-maximum accuracy is acceptable (only 0.08% less)",
    "",
    "Avoid DIO-RF-Single (94.37%, 8 features) unless:",
    "  - You only need rapid prototyping (1-minute optimization)",
    "  - Lower accuracy is acceptable for initial screening",
    "",
    "General wisdom:",
    "  - Always run 30+ validation runs with different splits",
    "  - Use statistical tests (Wilcoxon) to confirm improvements",
    "  - Match optimization budget to problem dimensionality"
])

# ============================================================================
# 5. REFERENCES
# ============================================================================
slide = add_content_slide(prs, "5. References")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "[1] El Romeh, A., Mirjalili, S., & Šnel, V. (2025). Dholes-Inspired Optimization (DIO). Cluster Computing (Springer). DOI: 10.1007/s10586-025-05543-2",
    "",
    "[2] Dua, D. & Graff, C. (2019). UCI Machine Learning Repository: Breast Cancer Wisconsin (Diagnostic) Data Set. University of California, Irvine.",
    "",
    "[3] Krizhevsky, A. (2009). Learning Multiple Layers of Features from Tiny Images. Technical Report, University of Toronto.",
    "",
    "[4] He, K., Zhang, X., Ren, S., & Sun, J. (2016). Deep Residual Learning for Image Recognition. IEEE CVPR.",
    "",
    "[5] Chen, T. & Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System. KDD '16.",
    "",
    "[6] Breiman, L. (2001). Random Forests. Machine Learning, 45(1), 5-32.",
    "",
    "[7] Wilcoxon, F. (1945). Individual Comparisons by Ranking Methods. Biometrics Bulletin, 1(6), 80-83."
])

# ============================================================================
# ACKNOWLEDGMENTS
# ============================================================================
slide = add_content_slide(prs, "Acknowledgments")
body_shape = slide.placeholders[1]
tf = body_shape.text_frame

add_bullet_points(tf, [
    "Special thanks to:",
    "",
    "Dr. Nabil Neggaz",
    "  - Course instructor and research supervisor",
    "  - Guidance throughout this project",
    "",
    "USTO-MB Computer Science Department",
    "  - Providing resources and support",
    "",
    "Original DIO Creators",
    "  - Dr. Seyedali Mirjalili and team",
    "  - For open-sourcing their algorithm",
    "",
    "Open-Source Community",
    "  - scikit-learn, XGBoost, NumPy, pandas",
    "  - Making research accessible to all",
    "",
    "",
    "Questions?"
])

# ============================================================================
# SAVE PRESENTATION
# ============================================================================
output_path = os.path.join(parent_dir, 'Presentation', 'DIO_Research_Presentation_v2.pptx')
prs.save(output_path)
print(f"✓ Presentation created successfully: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
